"""
보고 자동화 (Report Generator)
- 과거 보고자료 패턴 학습
- Claude API 연동 초안 생성
- PPTX 템플릿 기반 내보내기
"""
import json
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional
import re


class ReportTemplate:
    """보고자료 템플릿 분석 및 매칭"""

    def __init__(self, vault_dir: Path):
        self.vault_dir = vault_dir
        self.templates = {
            "weekly": [],      # 주간보고
            "monthly": [],     # 월간보고
            "cho": [],         # CHO 보고
            "team_intro": [],  # 팀소개
        }
        self._scan_templates()

    def _scan_templates(self):
        """템플릿 스캔"""
        for md_file in self.vault_dir.rglob("*.md"):
            name = md_file.name.lower()
            content = md_file.read_text(encoding='utf-8')[:1000]

            # 카테고리 분류
            if "주간" in name or "weekly" in name or "주차" in name:
                self.templates["weekly"].append({
                    "path": str(md_file),
                    "name": md_file.name,
                    "content": content
                })
            elif "cho" in name or "인재육성" in name:
                self.templates["cho"].append({
                    "path": str(md_file),
                    "name": md_file.name,
                    "content": content
                })
            elif "소개" in name or "intro" in name:
                self.templates["team_intro"].append({
                    "path": str(md_file),
                    "name": md_file.name,
                    "content": content
                })

        print(f"[TEMPLATE] 스캔 완료:")
        for category, docs in self.templates.items():
            print(f"  - {category}: {len(docs)}개")

    def get_template(self, report_type: str) -> Optional[Dict]:
        """최적 템플릿 반환"""
        docs = self.templates.get(report_type, [])
        if not docs:
            return None
        # 최신 문서 반환
        return max(docs, key=lambda x: x["path"])


class ReportGenerator:
    """Claude API 기반 보고 생성"""

    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key
        self.claude = None

        if api_key:
            try:
                import anthropic
                self.claude = anthropic.Anthropic(api_key=api_key)
                print("[INIT] Claude API 초기화 완료")
            except ImportError:
                print("[WARN] anthropic 패키지 없음: pip install anthropic")
            except Exception as e:
                print(f"[WARN] Claude API 초기화 실패: {e}")

    def generate_weekly_report(
        self,
        week: str,
        team: str,
        items: List[Dict],
        template: Optional[Dict] = None
    ) -> str:
        """
        주간보고 초안 생성

        Args:
            week: 주차 (예: "W48, 11 월 4 주차")
            team: 팀명 (예: "생산기술 HR 실")
            items: 완료 항목 목록 [{"task": "...", "result": "..."}]
            template: 템플릿 내용 (선택)
        """
        if not self.claude:
            # Claude 없이 로컬 생성
            return self._generate_local_weekly(week, team, items)

        # Claude API 사용
        prompt = self._build_weekly_prompt(week, team, items, template)

        response = self.claude.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=4096,
            messages=[
                {
                    "role": "user",
                    "content": prompt
                }
            ]
        )

        return response.content[0].text

    def _generate_local_weekly(
        self,
        week: str,
        team: str,
        items: List[Dict]
    ) -> str:
        """로컬 주간보고 생성 (Claude 없이)"""
        report = f"""# {team} 주간업무 보고 ({week})

## 1. 금주 완료 사항

"""
        for i, item in enumerate(items, 1):
            report += f"### {i}. {item.get('task', '작업명')}\n"
            report += f"- **진행상황**: {item.get('status', '완료')}\n"
            report += f"- **결과**: {item.get('result', '')}\n\n"

        report += """## 2. 차주 계획 사항

-
"""
        return report

    def _build_weekly_prompt(
        self,
        week: str,
        team: str,
        items: List[Dict],
        template: Optional[Dict] = None
    ) -> str:
        """Claude 프롬프트 구성"""
        prompt = f"""당신은 HR 팀의 전문 보고 작성 도우미입니다.

# 정보
- 팀: {team}
- 주차: {week}
- 완료 항목: {json.dumps(items, ensure_ascii=False)}

# 요구사항
1. 간결하고 명확한 비즈니스 문체 사용
2. 핵심 성과와 수치를 강조
3. 차주 계획은 구체적으로 작성
4. 한글 2 페이지 분량

# 출력 형식
마크다운 형식으로 작성하세요."""

        if template:
            prompt += f"\n\n# 참고 템플릿 스타일\n{template['content'][:500]}..."

        return prompt


class SlideExporter:
    """마크다운 → PPTX 내보내기"""

    def __init__(self):
        try:
            from pptx import Presentation
            self.Presentation = Presentation
            print("[INIT] python-pptx 사용 가능")
        except ImportError:
            print("[WARN] python-pptx 없음: pip install python-pptx")
            self.Presentation = None

    def export(
        self,
        content: str,
        output_path: Path,
        template_path: Optional[Path] = None
    ):
        """마크다운 내용을 PPTX 로 변환"""
        if not self.Presentation:
            print("[ERROR] PPTX 내보내기를 설치하려면: pip install python-pptx")
            return None

        if template_path and template_path.exists():
            prs = self.Presentation(template_path)
        else:
            prs = self.Presentation()

        # 마크다운 파싱
        slides = self._parse_markdown(content)

        for slide_data in slides:
            self._add_slide(prs, slide_data)

        prs.save(output_path)
        print(f"[EXPORT] PPTX 저장: {output_path}")
        return output_path

    def _parse_markdown(self, content: str) -> List[Dict]:
        """마크다운을 슬라이드 단위로 파싱"""
        slides = []
        current_slide = {"title": "", "content": []}

        for line in content.split('\n'):
            if line.startswith('# '):
                if current_slide["title"]:
                    slides.append(current_slide)
                current_slide = {"title": line[2:], "content": []}
            elif line.startswith('- '):
                current_slide["content"].append(line[2:])
            elif line.strip():
                current_slide["content"].append(line)

        if current_slide["title"]:
            slides.append(current_slide)

        return slides

    def _add_slide(self, prs, slide_data: Dict):
        """PPTX 슬라이 추가"""
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        # 제목
        slide.shapes.title.text = slide_data["title"]

        # 내용
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        for i, line in enumerate(slide_data["content"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.level = 0


def create_report_generator(
    vault_dir: Path,
    claude_api_key: Optional[str] = None
):
    """보고 생성기 팩토리"""
    template_mgr = ReportTemplate(vault_dir)
    report_gen = ReportGenerator(claude_api_key)
    slide_exp = SlideExporter()

    return {
        "templates": template_mgr,
        "generator": report_gen,
        "exporter": slide_exp
    }


if __name__ == "__main__":
    # 경로를 동적으로 계산
    BASE_DIR = Path(__file__).parent.parent

    # 보고 생성기 초기화
    tools = create_report_generator(BASE_DIR)

    # 데모: 주간보고 생성
    print("\n=== 보고 생성 데모 ===")

    sample_items = [
        {"task": "26 년 인재육성 업무계획 수립", "status": "완료", "result": "CHO 보고 완료"},
        {"task": "Domain AX 역량 인증제 운영", "status": "진행중", "result": "커미티 위원 구성 완료"},
        {"task": "Pulse 설문 분석", "status": "완료", "result": "EX Intelligence 리포트 작성"},
    ]

    weekly_report = tools["generator"].generate_weekly_report(
        week="W16, 4 월 2 주차",
        team="생산기술 HR 실",
        items=sample_items
    )

    print("\n[생성된 보고 초안]")
    print(weekly_report[:1000])

    # PPTX 내보내기 (선택)
    output_path = BASE_DIR / "harness" / "output" / "weekly_report_demo.pptx"
    output_path.parent.mkdir(parents=True, exist_ok=True)

    tools["exporter"].export(weekly_report, output_path)
